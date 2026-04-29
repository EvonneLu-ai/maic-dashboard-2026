#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
MAIC 2026 每週報告自動產出腳本
用法：python generate_report.py <csv檔案路徑> [上週csv路徑（可選）]

輸出：
  output/maic_report_YYYYMMDD.html   → HTML 互動報表
  output/maic_pdf_YYYYMMDD.html      → PDF 精簡報告（用瀏覽器列印）
  output/maic_slides_YYYYMMDD.pptx   → PPTX 投影片
  output/maic_summary_YYYYMMDD.txt   → 摘要文字訊息（可貼 LINE/Slack）
"""

import sys
import os
import csv
import json
from datetime import datetime
from collections import defaultdict, Counter

# ── 讀取 CSV ──────────────────────────────────────────────
def load_csv(path):
    rows = []
    with open(path, encoding='utf-8-sig') as f:
        reader = csv.DictReader(f)
        for row in reader:
            row['學生人數'] = int(row.get('學生人數', 0) or 0)
            row['指導老師數'] = int(row.get('指導老師數', 0) or 0)
            rows.append(row)
    return rows

# ── 統計分析 ──────────────────────────────────────────────
def analyze(rows, prev_rows=None):
    total_teams = len(rows)
    total_students = sum(r['學生人數'] for r in rows)
    total_teachers = sum(r['指導老師數'] for r in rows)
    schools = set(r['學校'] for r in rows)
    total_schools = len(schools)

    # 賽道分布
    track_count = Counter(r['賽道'] for r in rows)

    # 學校排行
    school_stats = defaultdict(lambda: {'teams': 0, 'students': 0, 'teachers': 0,
                                         '創意': 0, '創新': 0, '創業': 0})
    for r in rows:
        s = r['學校']
        school_stats[s]['teams'] += 1
        school_stats[s]['students'] += r['學生人數']
        school_stats[s]['teachers'] += r['指導老師數']
        school_stats[s][r['賽道']] += 1
    school_rank = sorted(school_stats.items(), key=lambda x: -x[1]['teams'])

    # 科系排行
    dept_count = Counter(r['科系'] for r in rows)
    dept_rank = dept_count.most_common(25)

    # 週增量（與上週比較）
    prev_teams = len(prev_rows) if prev_rows else None
    prev_schools_set = set(r['學校'] for r in prev_rows) if prev_rows else None
    delta_teams = (total_teams - prev_teams) if prev_teams is not None else None
    delta_schools = (total_schools - len(prev_schools_set)) if prev_schools_set is not None else None

    # 報名日期趨勢（按週累計）
    date_counter = Counter()
    for r in rows:
        d = r.get('報名日期', '')
        if d:
            date_counter[d[:7]] += 1  # 按月統計

    return {
        'total_teams': total_teams,
        'total_students': total_students,
        'total_teachers': total_teachers,
        'total_schools': total_schools,
        'track_count': dict(track_count),
        'school_rank': school_rank,
        'dept_rank': dept_rank,
        'delta_teams': delta_teams,
        'delta_schools': delta_schools,
        'date_trend': sorted(date_counter.items()),
    }

# ── ① HTML 互動報表 ──────────────────────────────────────
def gen_html(stats, date_str, out_path):
    school_rows_html = ''
    for i, (school, s) in enumerate(stats['school_rank'], 1):
        delta_html = '<span style="color:#8e8e93">—</span>'
        school_rows_html += f"""
        <tr>
          <td class="rank">{i}</td>
          <td class="school-name">{school}</td>
          <td class="num-cell">{s['teams']}</td>
          <td class="delta-cell">{delta_html}</td>
          <td><span class="track-badge c1">{s['創意']}</span></td>
          <td><span class="track-badge c2">{s['創新']}</span></td>
          <td><span class="track-badge c3">{s['創業']}</span></td>
          <td class="num-cell">{s['students']} 人<small style="color:#8e8e93"> ＋{s['teachers']} 師</small></td>
        </tr>"""

    dept_rows_html = ''
    for i, (dept, cnt) in enumerate(stats['dept_rank'], 1):
        dept_rows_html += f'<tr><td class="rank">{i}</td><td>{dept}</td><td class="num-cell">{cnt}</td></tr>'

    track = stats['track_count']
    t_total = stats['total_teams']
    html = f"""<!DOCTYPE html>
<html lang="zh-TW">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>MAIC 2026 週報 {date_str}</title>
<style>
*, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
body {{ font-family: -apple-system, 'PingFang TC', sans-serif; background: #f5f5f7; color: #1d1d1f; }}
.page {{ max-width: 1100px; margin: 0 auto; padding: 32px 24px 64px; }}
.header {{ background: linear-gradient(135deg, #0071e3, #34c759); border-radius: 20px; padding: 28px 32px; color: #fff; margin-bottom: 28px; }}
.header h1 {{ font-size: 26px; font-weight: 800; }}
.header p  {{ font-size: 13px; opacity: 0.85; margin-top: 4px; }}
.kpi-row {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 14px; margin-bottom: 24px; }}
.kpi {{ background: #fff; border-radius: 16px; padding: 20px; text-align: center; box-shadow: 0 2px 12px rgba(0,0,0,0.06); }}
.kpi .num {{ font-size: 36px; font-weight: 800; color: #0071e3; }}
.kpi .lbl {{ font-size: 12px; color: #8e8e93; margin-top: 4px; }}
.kpi .chg {{ font-size: 12px; color: #34c759; margin-top: 2px; font-weight: 600; }}
.track-row {{ display: flex; gap: 12px; margin-bottom: 24px; }}
.track-pill {{ flex: 1; border-radius: 14px; padding: 16px; text-align: center; color: #fff; font-weight: 700; font-size: 20px; }}
.t-c1 {{ background: linear-gradient(135deg, #ff9f0a, #ff6b00); }}
.t-c2 {{ background: linear-gradient(135deg, #34c759, #1a8c36); }}
.t-c3 {{ background: linear-gradient(135deg, #0071e3, #0050a0); }}
.track-pill .lbl {{ font-size: 12px; font-weight: 500; margin-top: 4px; opacity: 0.9; }}
.section {{ background: #fff; border-radius: 16px; padding: 24px; margin-bottom: 20px; box-shadow: 0 2px 12px rgba(0,0,0,0.06); }}
.section h2 {{ font-size: 16px; font-weight: 700; margin-bottom: 16px; color: #1d1d1f; }}
.search-bar {{ width: 100%; padding: 10px 14px; border: 1px solid #d1d1d6; border-radius: 10px; font-size: 14px; margin-bottom: 14px; outline: none; }}
.search-bar:focus {{ border-color: #0071e3; }}
table {{ width: 100%; border-collapse: collapse; font-size: 13px; }}
th {{ background: #f5f5f7; padding: 9px 10px; text-align: left; font-size: 11px; font-weight: 700; color: #8e8e93; letter-spacing: 0.5px; text-transform: uppercase; }}
td {{ padding: 10px 10px; border-bottom: 1px solid #f0f0f0; }}
tr:last-child td {{ border-bottom: none; }}
tr:hover td {{ background: #f9f9fb; }}
.rank {{ font-size: 12px; color: #8e8e93; font-weight: 700; width: 30px; }}
.num-cell {{ text-align: right; font-weight: 700; color: #1d1d1f; }}
.delta-cell {{ text-align: center; }}
.school-name {{ font-weight: 600; }}
.track-badge {{ display: inline-block; padding: 2px 8px; border-radius: 20px; font-size: 11px; font-weight: 700; }}
.c1 {{ background: rgba(255,159,10,0.15); color: #b85a00; }}
.c2 {{ background: rgba(52,199,89,0.15); color: #1a6628; }}
.c3 {{ background: rgba(0,113,227,0.15); color: #0050a0; }}
.footer {{ text-align: center; font-size: 11px; color: #8e8e93; margin-top: 32px; }}
@media (max-width: 600px) {{
  .kpi-row {{ grid-template-columns: repeat(2, 1fr); }}
  .track-row {{ flex-direction: column; }}
}}
</style>
</head>
<body>
<div class="page">

  <div class="header">
    <h1>📊 MAIC 2026 · 週報 {date_str}</h1>
    <p>行動應用創新賽 2026 ｜ 報名數據即時總覽</p>
  </div>

  <div class="kpi-row">
    <div class="kpi">
      <div class="num">{stats['total_teams']}</div>
      <div class="lbl">報名組數</div>
      <div class="chg">{f'▲ +{stats["delta_teams"]}' if stats['delta_teams'] else '—'}</div>
    </div>
    <div class="kpi">
      <div class="num">{stats['total_students']}</div>
      <div class="lbl">參賽學生人數</div>
      <div class="chg">含指導老師 {stats['total_teachers']} 人</div>
    </div>
    <div class="kpi">
      <div class="num">{stats['total_schools']}</div>
      <div class="lbl">參與學校數</div>
      <div class="chg">{f'▲ +{stats["delta_schools"]} 所' if stats['delta_schools'] else '—'}</div>
    </div>
    <div class="kpi">
      <div class="num">3</div>
      <div class="lbl">競賽賽道</div>
      <div class="chg">創意 / 創新 / 創業</div>
    </div>
  </div>

  <div class="track-row">
    <div class="track-pill t-c1">
      {track.get('創意', 0)}
      <div class="lbl">🎨 創意賽道</div>
    </div>
    <div class="track-pill t-c2">
      {track.get('創新', 0)}
      <div class="lbl">🔬 創新賽道</div>
    </div>
    <div class="track-pill t-c3">
      {track.get('創業', 0)}
      <div class="lbl">🚀 創業賽道</div>
    </div>
  </div>

  <div class="section">
    <h2>🏫 各校組數排行</h2>
    <input class="search-bar" id="school-search" placeholder="🔍 搜尋學校名稱..." oninput="filterTable('school-table', this.value, 1)">
    <table id="school-table">
      <thead><tr>
        <th>#</th><th>學校</th><th style="text-align:right">組數</th><th style="text-align:center">週比較</th>
        <th>創意</th><th>創新</th><th>創業</th><th style="text-align:right">人數</th>
      </tr></thead>
      <tbody>{school_rows_html}</tbody>
    </table>
  </div>

  <div class="section">
    <h2>📚 科系報名排行（前 25）</h2>
    <input class="search-bar" id="dept-search" placeholder="🔍 搜尋科系名稱..." oninput="filterTable('dept-table', this.value, 1)">
    <table id="dept-table">
      <thead><tr>
        <th>#</th><th>科系 / 所</th><th style="text-align:right">組數</th>
      </tr></thead>
      <tbody>{dept_rows_html}</tbody>
    </table>
  </div>

  <div class="footer">
    MAIC 行動應用創新賽 2026 · 週報 {date_str}<br>
    <a href="https://evonnelu-ai.github.io/maic-dashboard-2026/" style="color:#0071e3">→ 完整 Dashboard</a>
  </div>

</div>
<script>
function filterTable(tableId, query, colIdx) {{
  const q = query.toLowerCase();
  const rows = document.getElementById(tableId).querySelectorAll('tbody tr');
  rows.forEach(r => {{
    const cell = r.cells[colIdx];
    r.style.display = (!q || cell.textContent.toLowerCase().includes(q)) ? '' : 'none';
  }});
}}
</script>
</body>
</html>"""
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(html)
    print(f"  ✅ HTML 報表：{out_path}")

# ── ② PDF 精簡報告（print-friendly HTML）────────────────
def gen_pdf_html(stats, date_str, out_path):
    school_rows = ''
    for i, (school, s) in enumerate(stats['school_rank'][:10], 1):
        school_rows += f"<tr><td>{i}</td><td>{school}</td><td>{s['teams']}</td><td>{s['創意']}</td><td>{s['創新']}</td><td>{s['創業']}</td><td>{s['students']}人+{s['teachers']}師</td></tr>"

    dept_rows = ''
    for i, (dept, cnt) in enumerate(stats['dept_rank'][:10], 1):
        dept_rows += f"<tr><td>{i}</td><td>{dept}</td><td>{cnt}</td></tr>"

    track = stats['track_count']
    html = f"""<!DOCTYPE html>
<html lang="zh-TW">
<head>
<meta charset="UTF-8">
<title>MAIC 2026 PDF 週報 {date_str}</title>
<style>
@page {{ size: A4 portrait; margin: 18mm 15mm 15mm 15mm; }}
*, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
body {{ font-family: -apple-system, 'PingFang TC', 'Helvetica Neue', sans-serif; font-size: 11pt; color: #1d1d1f; background: #fff; }}
.header {{ border-bottom: 3px solid #0071e3; padding-bottom: 10px; margin-bottom: 16px; display: flex; justify-content: space-between; align-items: flex-end; }}
.header h1 {{ font-size: 18pt; font-weight: 800; color: #0071e3; }}
.header .meta {{ font-size: 9pt; color: #8e8e93; text-align: right; }}
.kpi-grid {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 10px; margin-bottom: 16px; }}
.kpi {{ border: 1.5px solid #d1d1d6; border-radius: 10px; padding: 10px; text-align: center; }}
.kpi .num {{ font-size: 22pt; font-weight: 800; color: #0071e3; }}
.kpi .lbl {{ font-size: 8pt; color: #8e8e93; }}
.track-grid {{ display: grid; grid-template-columns: repeat(3, 1fr); gap: 10px; margin-bottom: 16px; }}
.track {{ border-radius: 10px; padding: 10px; text-align: center; color: #fff; }}
.t1 {{ background: #ff9f0a; }}
.t2 {{ background: #34c759; }}
.t3 {{ background: #0071e3; }}
.track .cnt {{ font-size: 20pt; font-weight: 800; }}
.track .lbl {{ font-size: 9pt; opacity: 0.9; }}
.tables {{ display: grid; grid-template-columns: 1.6fr 1fr; gap: 16px; margin-bottom: 16px; }}
h2 {{ font-size: 11pt; font-weight: 700; margin-bottom: 8px; color: #1d1d1f; border-left: 3px solid #0071e3; padding-left: 8px; }}
table {{ width: 100%; border-collapse: collapse; font-size: 9pt; }}
th {{ background: #f5f5f7; padding: 5px 7px; text-align: left; font-weight: 700; color: #8e8e93; font-size: 8pt; }}
td {{ padding: 5px 7px; border-bottom: 1px solid #f0f0f0; }}
td:first-child {{ color: #8e8e93; font-weight: 700; width: 22px; }}
.footer {{ margin-top: 12px; text-align: center; font-size: 8pt; color: #8e8e93; border-top: 1px solid #e0e0e5; padding-top: 8px; }}
.print-btn {{ position: fixed; bottom: 24px; right: 24px; background: #0071e3; color: #fff; border: none; border-radius: 12px; padding: 12px 22px; font-size: 14px; font-weight: 700; cursor: pointer; box-shadow: 0 4px 16px rgba(0,113,227,0.35); }}
.print-btn:hover {{ background: #0060c7; }}
@media print {{ .print-btn {{ display: none; }} }}
</style>
</head>
<body>
<div class="header">
  <div><h1>📊 MAIC 2026 週報</h1><div style="font-size:9pt;color:#8e8e93;margin-top:3px">行動應用創新賽 · 報名數據摘要</div></div>
  <div class="meta">資料截至：{date_str}<br>Apple EDU Taiwan · Evonne Lu</div>
</div>

<div class="kpi-grid">
  <div class="kpi"><div class="num">{stats['total_teams']}</div><div class="lbl">報名組數</div></div>
  <div class="kpi"><div class="num">{stats['total_students']}</div><div class="lbl">學生人數</div></div>
  <div class="kpi"><div class="num">{stats['total_schools']}</div><div class="lbl">學校數</div></div>
  <div class="kpi"><div class="num">{stats['total_teams'] + stats['total_teachers']}</div><div class="lbl">總記錄筆數</div></div>
</div>

<div class="track-grid">
  <div class="track t1"><div class="cnt">{track.get('創意', 0)}</div><div class="lbl">🎨 創意賽道</div></div>
  <div class="track t2"><div class="cnt">{track.get('創新', 0)}</div><div class="lbl">🔬 創新賽道</div></div>
  <div class="track t3"><div class="cnt">{track.get('創業', 0)}</div><div class="lbl">🚀 創業賽道</div></div>
</div>

<div class="tables">
  <div>
    <h2>各校組數排行（Top 10）</h2>
    <table>
      <thead><tr><th>#</th><th>學校</th><th>組</th><th>意</th><th>新</th><th>業</th><th>人數</th></tr></thead>
      <tbody>{school_rows}</tbody>
    </table>
  </div>
  <div>
    <h2>科系排行（Top 10）</h2>
    <table>
      <thead><tr><th>#</th><th>科系</th><th>組</th></tr></thead>
      <tbody>{dept_rows}</tbody>
    </table>
  </div>
</div>

<div class="footer">
  MAIC 2026 週報 {date_str} ｜ Dashboard：https://evonnelu-ai.github.io/maic-dashboard-2026/
</div>
<button class="print-btn" onclick="window.print()">🖨️ 列印 / 儲存 PDF</button>
</body>
</html>"""
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(html)
    print(f"  ✅ PDF 頁面：{out_path}（用瀏覽器開啟後列印）")

# ── ③ PPTX 投影片 ─────────────────────────────────────────
def gen_pptx(stats, date_str, out_path):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("  ⚠️  python-pptx 未安裝，跳過 PPTX 輸出。")
        print("      請執行：pip install python-pptx")
        return

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BLUE  = RGBColor(0x00, 0x71, 0xE3)
    GREEN = RGBColor(0x34, 0xC7, 0x59)
    ORANGE= RGBColor(0xFF, 0x9F, 0x0A)
    DARK  = RGBColor(0x1d, 0x1d, 0x1f)
    GRAY  = RGBColor(0x8e, 0x8e, 0x93)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    blank = prs.slide_layouts[6]  # 全空白

    def add_textbox(slide, text, left, top, width, height,
                    size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return tb

    def add_rect(slide, left, top, width, height, fill_color, radius=None):
        from pptx.util import Inches
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    # ── Slide 1: 封面 ──
    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, 13.33, 7.5, RGBColor(0x00, 0x47, 0x9A))
    add_textbox(s1, 'MAIC 2026', 1.5, 1.8, 10, 1.2, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s1, f'每週報名數據 週報 · {date_str}', 1.5, 3.2, 10, 0.7, size=22, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
    add_textbox(s1, 'Apple EDU Taiwan · Evonne Lu', 1.5, 4.0, 10, 0.5, size=16, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

    # ── Slide 2: KPI 總覽 ──
    track = stats['track_count']
    s2 = prs.slides.add_slide(blank)
    add_rect(s2, 0, 0, 13.33, 1.1, BLUE)
    add_textbox(s2, f'📊 本週數據總覽', 0.3, 0.1, 12, 0.9, size=24, bold=True, color=WHITE)
    add_textbox(s2, f'資料截至 {date_str}', 10.0, 0.25, 3, 0.6, size=13, color=RGBColor(0xAA, 0xCC, 0xFF))

    kpis = [
        (str(stats['total_teams']), '報名組數', 0.5, BLUE),
        (str(stats['total_students']), '學生人數', 3.5, GREEN),
        (str(stats['total_schools']), '參與學校數', 6.5, ORANGE),
        (f"{track.get('創意',0)}/{track.get('創新',0)}/{track.get('創業',0)}", '意/新/業', 9.5, RGBColor(0xAF,0x52,0xDE)),
    ]
    for val, lbl, x, clr in kpis:
        add_rect(s2, x, 1.5, 2.8, 2.5, RGBColor(0xF5,0xF5,0xF7))
        add_textbox(s2, val, x, 1.7, 2.8, 1.2, size=40, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_textbox(s2, lbl, x, 3.0, 2.8, 0.5, size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # ── Slide 3: 學校排行 ──
    s3 = prs.slides.add_slide(blank)
    add_rect(s3, 0, 0, 13.33, 1.1, GREEN)
    add_textbox(s3, '🏫 各校組數排行（Top 10）', 0.3, 0.1, 12, 0.9, size=24, bold=True, color=WHITE)
    headers = ['#', '學校', '組數', '創意', '創新', '創業', '人數']
    widths  = [0.4, 3.5, 0.8, 0.8, 0.8, 0.8, 1.5]
    x_pos   = [0.3, 0.8, 4.4, 5.3, 6.2, 7.1, 8.0]
    for j, (h, xp) in enumerate(zip(headers, x_pos)):
        add_textbox(s3, h, xp, 1.3, widths[j], 0.4, size=11, bold=True, color=GRAY)
    for i, (school, s) in enumerate(stats['school_rank'][:10]):
        y = 1.8 + i * 0.48
        if i % 2 == 0:
            add_rect(s3, 0.3, y - 0.05, 9.5, 0.42, RGBColor(0xF9,0xF9,0xFB))
        row_vals = [str(i+1), school, str(s['teams']), str(s['創意']), str(s['創新']), str(s['創業']), f"{s['students']}人"]
        for j, (val, xp, w) in enumerate(zip(row_vals, x_pos, widths)):
            sz = 13 if j != 1 else 12
            clr = BLUE if j == 2 else DARK
            add_textbox(s3, val, xp, y, w, 0.4, size=sz, color=clr, bold=(j==2))

    # ── Slide 4: 科系 + 賽道 ──
    s4 = prs.slides.add_slide(blank)
    add_rect(s4, 0, 0, 13.33, 1.1, ORANGE)
    add_textbox(s4, '📚 科系分布 & 賽道佔比', 0.3, 0.1, 12, 0.9, size=24, bold=True, color=WHITE)
    add_textbox(s4, '科系排行（Top 10）', 0.5, 1.3, 5, 0.4, size=14, bold=True, color=DARK)
    dept_x = [0.5, 1.2, 4.8]
    dept_w = [0.5, 3.5, 0.8]
    for i, (dept, cnt) in enumerate(stats['dept_rank'][:10]):
        y = 1.8 + i * 0.46
        if i % 2 == 0:
            add_rect(s4, 0.5, y - 0.05, 5.0, 0.4, RGBColor(0xF9,0xF9,0xFB))
        for val, xp, w in zip([str(i+1), dept, str(cnt)], dept_x, dept_w):
            add_textbox(s4, val, xp, y, w, 0.38, size=12, color=DARK if xp != 4.8 else ORANGE, bold=(xp==4.8))

    add_textbox(s4, '賽道分布', 7.2, 1.3, 5, 0.4, size=14, bold=True, color=DARK)
    tracks = [('🎨 創意', track.get('創意',0), ORANGE), ('🔬 創新', track.get('創新',0), GREEN), ('🚀 創業', track.get('創業',0), BLUE)]
    total_t = sum(t[1] for t in tracks) or 1
    for i, (name, cnt, clr) in enumerate(tracks):
        y = 2.0 + i * 1.3
        add_rect(s4, 7.2, y, 5.0, 1.0, RGBColor(0xF5,0xF5,0xF7))
        bar_w = max(0.2, 4.5 * cnt / total_t)
        add_rect(s4, 7.2, y, bar_w, 1.0, clr)
        add_textbox(s4, f'{name}  {cnt} 組 ({cnt*100//total_t}%)', 7.4, y + 0.25, 4.5, 0.5, size=16, bold=True, color=WHITE if bar_w > 1 else DARK)

    prs.save(out_path)
    print(f"  ✅ PPTX 投影片：{out_path}")

# ── ④ 摘要文字訊息 ───────────────────────────────────────
def gen_summary(stats, date_str, out_path):
    track = stats['track_count']
    top3 = [f"{s}（{d['teams']}組）" for s, d in stats['school_rank'][:3]]
    delta_str = f"（較上週 ▲+{stats['delta_teams']}）" if stats['delta_teams'] else ''
    summary = f"""📊 MAIC 2026 週報 {date_str}
━━━━━━━━━━━━━━━━━━━━
🏆 報名組數：{stats['total_teams']} 組 {delta_str}
👥 學生人數：{stats['total_students']} 人（含 {stats['total_teachers']} 位指導老師）
🏫 涵蓋學校：{stats['total_schools']} 所
🎯 賽道分布：創意 {track.get('創意',0)} ｜ 創新 {track.get('創新',0)} ｜ 創業 {track.get('創業',0)}

🔝 學校 Top 3：
   {' / '.join(top3)}

🔗 完整 Dashboard：
   https://evonnelu-ai.github.io/maic-dashboard-2026/
━━━━━━━━━━━━━━━━━━━━
Apple EDU Taiwan · Evonne Lu"""
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(summary)
    print(f"  ✅ 摘要訊息：{out_path}")
    print()
    print("── 訊息預覽 ──────────────────────────────────")
    print(summary)
    print("──────────────────────────────────────────────")

# ── 主程式 ────────────────────────────────────────────────
def main():
    if len(sys.argv) < 2:
        print("用法：python generate_report.py <csv檔案路徑> [上週csv路徑]")
        print("範例：python generate_report.py sample_data.csv")
        sys.exit(1)

    csv_path = sys.argv[1]
    prev_csv = sys.argv[2] if len(sys.argv) > 2 else None

    if not os.path.exists(csv_path):
        print(f"❌ 找不到檔案：{csv_path}")
        sys.exit(1)

    rows = load_csv(csv_path)
    prev_rows = load_csv(prev_csv) if prev_csv else None
    stats = analyze(rows, prev_rows)

    date_str = datetime.today().strftime('%Y-%m-%d')
    out_dir = os.path.join(os.path.dirname(csv_path), 'output')
    os.makedirs(out_dir, exist_ok=True)

    print(f"\n🚀 MAIC 2026 週報產出中... ({date_str})")
    print(f"   資料：{len(rows)} 筆 | 學校：{stats['total_schools']} 所 | 組數：{stats['total_teams']}\n")

    gen_html(stats, date_str,
             os.path.join(out_dir, f'maic_report_{date_str.replace("-","")}.html'))
    gen_pdf_html(stats, date_str,
             os.path.join(out_dir, f'maic_pdf_{date_str.replace("-","")}.html'))
    gen_pptx(stats, date_str,
             os.path.join(out_dir, f'maic_slides_{date_str.replace("-","")}.pptx'))
    gen_summary(stats, date_str,
             os.path.join(out_dir, f'maic_summary_{date_str.replace("-","")}.txt'))

    print(f"\n✨ 全部完成！輸出目錄：{out_dir}/")

if __name__ == '__main__':
    main()
