#!/usr/bin/env python3
"""
MAIC Dashboard 一鍵四出
════════════════════════
cd ~/Desktop/maic-dashboard-deploy/maic_report
python3 generate_from_dashboard.py

輸出：
  ① MAIC_Dashboard_YYYY-MM-DD.html   — 互動報表快照
  ② MAIC_Report_YYYY-MM-DD.pdf       — A4 精簡報告
  ③ MAIC_Slides_YYYY-MM-DD.pptx      — 週會投影片
  ④ MAIC_Summary_YYYY-MM-DD.txt      — LINE/Slack 摘要訊息
"""

import re, os, sys, shutil, subprocess
from datetime import datetime
from pathlib import Path

# ── 本地套件路徑（lib/ 目錄，免系統安裝）─────────────────────────
_LIB = Path(__file__).parent / "lib"
if _LIB.exists() and str(_LIB) not in sys.path:
    sys.path.insert(0, str(_LIB))

# ── 路徑設定 ──────────────────────────────────────────────────────
BASE_DIR    = Path(__file__).parent.parent          # maic-dashboard-deploy/
SOURCE_HTML = BASE_DIR / "index.html"
OUTPUT_DIR  = Path(__file__).parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)
TODAY       = datetime.now().strftime("%Y-%m-%d")
DASHBOARD_URL = "https://evonnelu-ai.github.io/maic-dashboard-2026/"

# ── 讀取 index.html ───────────────────────────────────────────────
if not SOURCE_HTML.exists():
    print(f"❌ 找不到 {SOURCE_HTML}，請確認路徑"); sys.exit(1)
html = SOURCE_HTML.read_text(encoding="utf-8")

# ── 解析函式 ──────────────────────────────────────────────────────
def kpi(label):
    m = re.search(rf"label:\s*'{re.escape(label)}'[^}}]*?value:\s*'([^']+)'", html)
    return m.group(1).strip() if m else "—"

def kpi_sub(label):
    m = re.search(rf"label:\s*'{re.escape(label)}'[^}}]*?sub:\s*'([^']+)'", html)
    return m.group(1).strip() if m else ""

def meta(key):
    m = re.search(rf"{re.escape(key)}:\s*'([^']+)'", html)
    return m.group(1).strip() if m else ""

def group_val(key):
    m = re.search(rf"{re.escape(key)}:\s*(\d+)", html)
    return m.group(1) if m else "—"

# ── 關鍵數據 ──────────────────────────────────────────────────────
date_raw   = meta("dataAsOf")
date_label = re.sub(r"[📅\s]*資料截至\s*", "", date_raw).strip() or TODAY

teams      = kpi("累計報名隊伍")
rate       = kpi("目標達成率")
accounts   = kpi("累計註冊帳號")
students   = kpi("報名學生人數")
schools    = kpi("涵蓋學校數")
sessions_n = kpi("已舉辦場次")
attend     = kpi("宣講累計出席")
next_sess  = kpi("下一場（預計）")

sub_teams  = kpi_sub("累計報名隊伍")
sub_rate   = kpi_sub("目標達成率")
sub_sess   = kpi_sub("已舉辦場次")
sub_next   = kpi_sub("下一場（預計）")

m_new    = re.search(r"\+(\d+)\s*隊", sub_teams)
week_new = m_new.group(1) if m_new else "0"

m_goal   = re.search(r"（(\d+/\d+)）", sub_rate)
goal_str = m_goal.group(1) if m_goal else f"{teams}/800"

# 創意/創新/創業
creative   = group_val("creative")
innovation = group_val("innovative")
startup    = group_val("startup")

print(f"\n🚀 MAIC 一鍵四出 ｜ 資料截至 {date_label}")
print("=" * 50)


# ══════════════════════════════════════════════════════════
# 1️⃣  HTML — Dashboard 互動報表快照
# ══════════════════════════════════════════════════════════
html_out = OUTPUT_DIR / f"MAIC_Dashboard_{TODAY}.html"
shutil.copy2(SOURCE_HTML, html_out)
print(f"✅ HTML  → {html_out.name}")


# ══════════════════════════════════════════════════════════
# 2️⃣  PDF — A4 精簡報告（Chrome headless）
# ══════════════════════════════════════════════════════════
# ══════════════════════════════════════════════════════════
# 2️⃣  PDF — reportlab A4 精簡報告
# ══════════════════════════════════════════════════════════
pdf_out = OUTPUT_DIR / f"MAIC_Report_{TODAY}.pdf"

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import mm
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    Table, TableStyle, HRFlowable)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    # 嘗試載入中文字型
    _font = "Helvetica"
    for fp in ["/System/Library/Fonts/PingFang.ttc",
               "/Library/Fonts/Arial Unicode.ttf"]:
        if os.path.exists(fp):
            try:
                pdfmetrics.registerFont(TTFont("CJK", fp))
                _font = "CJK"
            except:
                pass
            break

    BLUE   = colors.HexColor("#0071e3")
    ORANGE = colors.HexColor("#ff9500")
    GREEN  = colors.HexColor("#34c759")
    PURPLE = colors.HexColor("#af52de")
    LIGHT  = colors.HexColor("#f5f5f7")
    DARK   = colors.HexColor("#1d1d1f")
    GRAY   = colors.HexColor("#86868b")

    doc = SimpleDocTemplate(str(pdf_out), pagesize=A4,
                            leftMargin=15*mm, rightMargin=15*mm,
                            topMargin=12*mm, bottomMargin=12*mm)
    W = A4[0] - 30*mm

    def ps(name, size=10, bold=False, color=DARK, align=TA_LEFT, leading=None):
        return ParagraphStyle(name, fontName=_font,
                              fontSize=size, textColor=color,
                              alignment=align,
                              leading=leading or size*1.35,
                              spaceAfter=2)

    story = []

    # ── Header ──────────────────────────────────────────
    hdr_data = [[
        Paragraph(f"<b>📱 MAIC 2026 週報</b>", ps("h1",18,color=colors.white,align=TA_LEFT)),
        Paragraph(f"📅 資料截至 {date_label}<br/>產出：{TODAY}",
                  ps("hd",9,color=colors.HexColor("#aaddff"),align=TA_RIGHT))
    ]]
    hdr = Table(hdr_data, colWidths=[W*0.65, W*0.35])
    hdr.setStyle(TableStyle([
        ("BACKGROUND", (0,0),(-1,-1), BLUE),
        ("ROUNDEDCORNERS", [6,6,6,6]),
        ("VALIGN",(0,0),(-1,-1),"MIDDLE"),
        ("TOPPADDING",(0,0),(-1,-1),8),
        ("BOTTOMPADDING",(0,0),(-1,-1),8),
        ("LEFTPADDING",(0,0),(0,-1),12),
        ("RIGHTPADDING",(-1,0),(-1,-1),12),
    ]))
    story += [hdr, Spacer(1, 6*mm)]

    # ── 區塊標題 ────────────────────────────────────────
    def section(title, color=BLUE):
        story.append(Table([[Paragraph(f"<b>{title}</b>",
                             ps("st",12,color=colors.white))]],
                           colWidths=[W],
                           style=[("BACKGROUND",(0,0),(-1,-1),color),
                                  ("TOPPADDING",(0,0),(-1,-1),5),
                                  ("BOTTOMPADDING",(0,0),(-1,-1),5),
                                  ("LEFTPADDING",(0,0),(-1,-1),8)]))
        story.append(Spacer(1,3*mm))

    # ── KPI 格子 ────────────────────────────────────────
    def kpi_table(items, cols=5):
        col_w = W / cols
        row = []
        for label, value, unit, color, sub in items:
            cell = Table([[Paragraph(f"<b>{label}</b>", ps("kl",8,color=color))],
                          [Paragraph(f"<b>{value}</b><font size='8'> {unit}</font>",
                                     ps("kv",18,color=DARK))],
                          [Paragraph(sub, ps("ks",7,color=GRAY))]],
                         colWidths=[col_w-6*mm])
            cell.setStyle(TableStyle([
                ("BACKGROUND",(0,0),(-1,-1),LIGHT),
                ("LINEABOVE",(0,0),(-1,0),3,color),
                ("TOPPADDING",(0,0),(-1,-1),4),
                ("BOTTOMPADDING",(0,0),(-1,-1),4),
                ("LEFTPADDING",(0,0),(-1,-1),5),
                ("RIGHTPADDING",(0,0),(-1,-1),5),
            ]))
            row.append(cell)
        while len(row) % cols != 0:
            row.append(Spacer(1,1))
        tbl = Table([row], colWidths=[col_w]*cols)
        tbl.setStyle(TableStyle([("VALIGN",(0,0),(-1,-1),"TOP")]))
        story.append(tbl)
        story.append(Spacer(1,4*mm))

    # ── 報名數據 ─────────────────────────────────────────
    section("📊 報名數據 KPI")
    kpi_table([
        ("累計報名隊伍", teams,      "組", BLUE,   f"本週新增 +{week_new} 隊"),
        ("目標達成率",   rate,       "",   ORANGE, f"{goal_str} / 目標 800"),
        ("累計註冊帳號", accounts,   "位", colors.HexColor("#5ac8fa"), "已驗證帳號"),
        ("報名學生人數", students,   "人", PURPLE, ""),
        ("涵蓋學校數",   schools,    "所", BLUE,   "全台各地"),
    ])

    # 進度條（文字版）
    try:
        pct_v = float(re.sub(r"[^\d.]","",rate))
    except:
        pct_v = 0
    bar_filled = int(pct_v / 100 * 30)
    bar_str = "█"*bar_filled + "░"*(30-bar_filled)
    story.append(Paragraph(
        f"目標進度　{bar_str}　{rate}（{goal_str} / 800）",
        ps("bar", 8, color=GRAY)))
    story.append(Spacer(1, 4*mm))

    # 賽道
    section("🏁 賽道分佈")
    kpi_table([
        ("🎨 創意組", creative,   "組", BLUE,   ""),
        ("💡 創新組", innovation, "組", ORANGE, ""),
        ("🚀 創業組", startup,    "組", GREEN,  ""),
    ], cols=3)

    # ── 宣講會 ───────────────────────────────────────────
    section("📣 宣講會成效", ORANGE)
    kpi_table([
        ("已舉辦場次",   sessions_n, "場", ORANGE, sub_sess),
        ("宣講累計出席", attend,     "人", GREEN,  "含課程公告 520 人"),
        ("下一場（預計）",next_sess, "",   PURPLE, sub_next),
    ], cols=3)

    # 場次表
    tbl_data = [["日期","學校","類型","時間","人數","場地"]]
    tbl_data += [
        ["3/18","高雄科技大學","線下","14:00–14:30","150","燕巢校區"],
        ["3/19","臺大（工作坊）","線下","09:30–12:30","15","共同教學館 103"],
        ["3/20","文化大學","線下","10:00–11:00","30","資訊傳播系"],
        ["3/25","海洋大學","線下","12:20–12:40","15","資訊工程學系"],
        ["4/1","慈濟大學","線下","16:10–16:30","45","經管系自媒體課"],
        ["4/2+9","東華大學","線上","12:30–13:30","10","線上"],
        ["4/16","東華大學資管系","公告","—","520","必修 R 語言課程"],
        ["4/24","逢甲大學","線下","15:40–17:00","20","資電 B29"],
    ]
    col_ws = [W*x for x in [0.08,0.22,0.09,0.18,0.08,0.35]]
    tbl = Table([[Paragraph(str(c), ps("tc",8,
                  color=(colors.white if ri==0 else DARK),
                  align=TA_CENTER)) for c in row]
                 for ri, row in enumerate(tbl_data)],
                colWidths=col_ws)
    ts = [("BACKGROUND",(0,0),(-1,0),ORANGE),
          ("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,LIGHT]),
          ("GRID",(0,0),(-1,-1),0.3,colors.HexColor("#e5e5ea")),
          ("TOPPADDING",(0,0),(-1,-1),3),("BOTTOMPADDING",(0,0),(-1,-1),3),
          ("LEFTPADDING",(0,0),(-1,-1),4),("RIGHTPADDING",(0,0),(-1,-1),4),
          ("VALIGN",(0,0),(-1,-1),"MIDDLE")]
    tbl.setStyle(TableStyle(ts))
    story += [tbl, Spacer(1,5*mm)]

    # ── Footer ───────────────────────────────────────────
    story.append(HRFlowable(width=W, color=colors.HexColor("#e5e5ea")))
    story.append(Spacer(1,2*mm))
    story.append(Table([[
        Paragraph(f"🔗 {DASHBOARD_URL}", ps("fu",8,color=BLUE)),
        Paragraph(f"MAIC 2026 · Straight A × Apple · {TODAY}",
                  ps("fr",8,color=GRAY,align=TA_RIGHT))
    ]], colWidths=[W*0.6, W*0.4]))

    doc.build(story)
    print(f"✅ PDF   → {pdf_out.name}")

except ImportError:
    print("⚠️  PDF  → reportlab 未安裝：pip3 install reportlab --target lib/")
except Exception as e:
    print(f"⚠️  PDF  → 發生錯誤：{e}")
    import traceback; traceback.print_exc()


# ══════════════════════════════════════════════════════════
# 3️⃣  PPTX — 週會投影片（python-pptx）
# ══════════════════════════════════════════════════════════
pptx_out = OUTPUT_DIR / f"MAIC_Slides_{TODAY}.pptx"

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # 色盤
    C_BLUE   = RGBColor(0x00, 0x71, 0xe3)
    C_ORANGE = RGBColor(0xff, 0x95, 0x00)
    C_GREEN  = RGBColor(0x34, 0xc7, 0x59)
    C_PURPLE = RGBColor(0xaf, 0x52, 0xde)
    C_WHITE  = RGBColor(0xff, 0xff, 0xff)
    C_DARK   = RGBColor(0x1d, 0x1d, 0x1f)
    C_LIGHT  = RGBColor(0xf5, 0xf5, 0xf7)
    C_GRAY   = RGBColor(0x86, 0x86, 0x8b)
    C_NAVY   = RGBColor(0x00, 0x4f, 0x9e)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    # ── 基礎繪圖函式 ─────────────────────────────
    def rect(slide, x, y, w, h, fill, line=False):
        s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = fill
        if not line: s.line.fill.background()
        return s

    def txt(slide, text, x, y, w, h, size=14, bold=False,
            color=None, align=PP_ALIGN.LEFT, italic=False, wrap=True):
        color = color or C_DARK
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.italic = italic
        return tb

    def kpi_card(slide, x, y, w, h, label, value, unit, color, sub=""):
        rect(slide, x, y, w, h, C_LIGHT)
        rect(slide, x, y, w, 0.07, color)
        txt(slide, label, x+0.12, y+0.13, w-0.24, 0.32, size=11, bold=True, color=color)
        txt(slide, value+unit, x+0.12, y+0.44, w-0.24, 0.72, size=26, bold=True, color=C_DARK)
        if sub:
            txt(slide, sub, x+0.12, y+h-0.42, w-0.24, 0.38, size=9, color=C_GRAY)

    def section_header(slide, title, color):
        rect(slide, 0, 0, 13.33, 1.05, color)
        txt(slide, title, 0.45, 0.18, 9, 0.72, size=30, bold=True, color=C_WHITE)
        txt(slide, f"資料截至 {date_label}", 9.5, 0.33, 3.6, 0.42,
            size=12, color=RGBColor(0xcc,0xdd,0xee), align=PP_ALIGN.RIGHT)

    # ── Slide 1：封面 ─────────────────────────────
    s1 = prs.slides.add_slide(BLANK)
    rect(s1, 0, 0, 13.33, 7.5, C_BLUE)
    rect(s1, 0, 0, 13.33, 0.25, C_NAVY)
    rect(s1, 0, 6.0, 13.33, 1.5, C_NAVY)
    # 裝飾圓
    for cx, cy, sz, op in [(11.5,1.0,2.5,C_NAVY),(12.5,3.0,1.5,C_NAVY),(10.8,5.0,1.2,C_NAVY)]:
        rect(s1, cx, cy, sz, sz, op)

    txt(s1, "MAIC 2026", 0, 1.5, 13.33, 2.0,
        size=72, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s1, "行動應用創新賽 · 每週報告", 0, 3.6, 13.33, 0.8,
        size=24, color=RGBColor(0xaa,0xd4,0xff), align=PP_ALIGN.CENTER)
    txt(s1, f"📅  資料截至 {date_label}", 0, 4.5, 13.33, 0.6,
        size=18, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s1, "Straight A  ×  Apple", 0, 6.2, 13.33, 0.6,
        size=14, color=RGBColor(0xaa,0xd4,0xff), align=PP_ALIGN.CENTER)

    # ── Slide 2：報名數據 KPI ─────────────────────
    s2 = prs.slides.add_slide(BLANK)
    section_header(s2, "📊  報名數據", C_BLUE)
    kpi_card(s2, 0.25, 1.2, 3.1, 2.0, "累計報名隊伍", teams, " 組", C_BLUE, f"本週新增 +{week_new} 隊")
    kpi_card(s2, 3.55, 1.2, 3.1, 2.0, "目標達成率", rate, "", C_ORANGE, f"{goal_str}  /  目標 800 組")
    kpi_card(s2, 6.85, 1.2, 3.1, 2.0, "累計註冊帳號", accounts, " 位", RGBColor(0x5a,0xc8,0xfa), "已驗證帳號")
    kpi_card(s2, 10.0, 1.2, 3.08, 2.0, "涵蓋學校數", schools, " 所", C_PURPLE, "全台各地")

    # 進度條
    try:
        pct_float = float(re.sub(r"[^\d.]", "", rate)) / 100
    except:
        pct_float = 0
    rect(s2, 0.25, 3.45, 12.83, 0.38, RGBColor(0xe5,0xe5,0xea))
    bar_w = 12.83 * pct_float
    if bar_w > 0:
        rect(s2, 0.25, 3.45, bar_w, 0.38, C_BLUE)
    txt(s2, f"目標進度：{teams} / 800 組  ({rate})", 0.25, 3.95, 12.83, 0.45,
        size=12, color=C_GRAY)

    # 賽道分佈
    rect(s2, 0.25, 4.55, 12.83, 0.38, RGBColor(0xe8,0xf4,0xfd))
    txt(s2, "🏁  賽道分佈", 0.45, 4.6, 4, 0.32, size=12, bold=True, color=C_BLUE)
    for i, (icon, label, val) in enumerate([
        ("🎨","創意組", creative), ("💡","創新組", innovation), ("🚀","創業組", startup)
    ]):
        kpi_card(s2, 0.25 + i*4.3, 5.1, 4.1, 1.9,
                 f"{icon} {label}", val, " 組",
                 [C_BLUE, C_ORANGE, C_GREEN][i])

    # ── Slide 3：行銷宣講 ─────────────────────────
    s3 = prs.slides.add_slide(BLANK)
    section_header(s3, "📣  宣講會成效", C_ORANGE)
    kpi_card(s3, 0.25, 1.2, 4.2, 2.0, "已舉辦場次", sessions_n, " 場", C_ORANGE, sub_sess)
    kpi_card(s3, 4.65, 1.2, 4.2, 2.0, "宣講累計出席", attend, " 人", C_GREEN, "含課程公告 520 人")
    kpi_card(s3, 9.1,  1.2, 4.0, 2.0, "最大單場人數", "150+", "", C_BLUE, "高科大 3/18")

    # 場次表格
    rect(s3, 0.25, 3.4, 12.83, 0.42, C_ORANGE)
    for ci, (col_txt, cx) in enumerate([("日期",0.3),("學校",1.2),("類型",5.5),("時間",6.7),("人數",8.5),("場地",9.5)]):
        txt(s3, col_txt, cx, 3.45, 1.5, 0.35, size=10, bold=True, color=C_WHITE)

    rows = [
        ("3/18","高雄科技大學","線下","14:00–14:30","150","燕巢校區"),
        ("3/19","臺灣大學（工作坊）","線下","09:30–12:30","15","共同教學館 103"),
        ("3/20","文化大學","線下","10:00–11:00","30","資訊傳播系"),
        ("3/25","海洋大學","線下","12:20–12:40","15","資訊工程學系"),
        ("4/1","慈濟大學","線下","16:10–16:30","45","經管系自媒體課"),
        ("4/2+9","東華大學","線上","12:30–13:30","10","線上"),
        ("4/16","東華大學資管系","公告","—","520","必修 R 語言課程"),
        ("4/24","逢甲大學","線下","15:40–17:00","20","資電 B29"),
    ]
    for ri, row in enumerate(rows):
        bg = C_LIGHT if ri % 2 == 0 else C_WHITE
        rect(s3, 0.25, 3.88+ri*0.37, 12.83, 0.35, bg)
        for ci, (val, cx) in enumerate(zip(row, [0.3,1.2,5.5,6.7,8.5,9.5])):
            txt(s3, val, cx, 3.9+ri*0.37, 1.5, 0.32, size=9, color=C_DARK)

    # ── Slide 4：摘要 & 下一步 ────────────────────
    s4 = prs.slides.add_slide(BLANK)
    section_header(s4, "📋  本週摘要 & 下一步", C_DARK)

    bullets = [
        (C_BLUE,   "🏁  累計報名",  f"報名隊伍 {teams} 組，本週新增 +{week_new} 隊，目標達成率 {rate}（{goal_str}）"),
        (C_ORANGE, "📣  宣講會",    f"已舉辦 {sessions_n} 場，宣講累計出席 {attend} 人，下一場：{next_sess}"),
        (C_GREEN,  "🏫  學校覆蓋",  f"已觸及 {schools} 所學校，報名學生 {students} 人"),
        (C_PURPLE, "⏭️  下一步",    "持續衝刺報名，密切追蹤企畫書上傳率，準備初審評分機制"),
        (C_BLUE,   "🎯  目標",      "衝刺目標 800 組，重點關注南部學校與企管/設計科系滲透率"),
    ]
    for i, (color, label, content) in enumerate(bullets):
        y = 1.3 + i * 1.08
        rect(s4, 0.25, y, 0.58, 0.72, color)
        txt(s4, label, 0.25, y, 0.58, 0.72,
            size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        rect(s4, 1.0, y, 12.1, 0.72, C_LIGHT)
        txt(s4, content, 1.15, y+0.08, 11.8, 0.56, size=14, color=C_DARK)

    # URL footer
    txt(s4, f"🔗  {DASHBOARD_URL}", 0.25, 6.9, 12.83, 0.45,
        size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

    prs.save(pptx_out)
    print(f"✅ PPTX  → {pptx_out.name}")

except ImportError:
    print("⚠️  PPTX → 請先安裝：pip3 install python-pptx")
except Exception as e:
    print(f"⚠️  PPTX → 發生錯誤：{e}")
    import traceback; traceback.print_exc()


# ══════════════════════════════════════════════════════════
# 4️⃣  摘要文字訊息 — LINE / Slack / Teams
# ══════════════════════════════════════════════════════════
txt_out = OUTPUT_DIR / f"MAIC_Summary_{TODAY}.txt"

summary = f"""📊 MAIC 2026 週報｜{date_label}
{'━' * 36}

🏁 累計報名隊伍　{teams} 組（本週 +{week_new} 隊）
🎯 目標達成率　　{rate}（{goal_str} ／目標 800 組）
👥 報名學生人數　{students} 人
🏫 涵蓋學校數　　{schools} 所

🎨 創意組 {creative} 組　💡 創新組 {innovation} 組　🚀 創業組 {startup} 組

{'━' * 36}

📣 宣講會已舉辦　{sessions_n} 場（{sub_sess}）
👥 宣講累計出席　{attend} 人
⏭️  下一場　　　　{next_sess}（{sub_next}）

{'━' * 36}

🔗 互動報表：{DASHBOARD_URL}

#MAIC2026 #行動應用創新賽 #StraightA #Apple
"""

txt_out.write_text(summary, encoding="utf-8")
print(f"✅ 摘要  → {txt_out.name}")


# ── 完成摘要 ────────────────────────────────────────────────
print()
print("=" * 50)
print(f"📁 輸出資料夾：{OUTPUT_DIR}")
print()
print(summary.strip())
print()

# 自動開啟輸出資料夾
subprocess.run(["open", str(OUTPUT_DIR)], check=False)
